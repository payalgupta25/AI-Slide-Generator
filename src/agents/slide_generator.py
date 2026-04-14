import os
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv


class SlideGenerator:
    def __init__(self, output_dir="output/"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def add_slide(self, title, content):
        prs=Presentation()
        
        #Title
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        #Content
        for i in range(0,len(content),5):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = f"{title} (part {i//5+1})"
            content_shape = slide.placeholders[1]
            
            for b in content[i:i+5]:
                p = content_shape.text_frame.add_paragraph()
                p.text = b

        # Save the presentation
        pptx_path = os.path.join(self.output_dir, "generated_presentation.pptx")
        prs.save(pptx_path)
        print(f"Presentation saved to {pptx_path}")

if __name__ == "__main__":
    sample_title= "Poverty"
    sample_points=  [
        "A1",
        "A2",
        "A3",
        "A4",
        "Its",
    ]

    slide_generator = SlideGenerator()
    slide_generator.add_slide(sample_title, sample_points)